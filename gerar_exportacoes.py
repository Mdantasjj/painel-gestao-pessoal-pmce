from pathlib import Path

from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION, XL_LEGEND_POSITION, XL_TICK_LABEL_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent
LOGO_IMAGE = ROOT / "assets" / "timbrado.png"
EXPORTS = ROOT / "exportacoes"
PDF_PATH = EXPORTS / "Painel_Gestao_Pessoal_PMCE.pdf"
PPTX_PATH = EXPORTS / "Painel_Gestao_Pessoal_PMCE_EDITAVEL.pptx"

GREEN_900 = "0D3C2D"
GREEN_800 = "13523B"
GREEN_600 = "1B8258"
GREEN_050 = "EFF8F3"
INK = "16231D"
MUTED = "6D7973"
LINE = "DFE6E2"
GOLD = "C1A253"
BLUE = "3B7E9D"
TEAL = "2B8982"
OLIVE = "698342"
WHITE = "FFFFFF"
BG = "F3F7F5"


def color(value: str) -> RGBColor:
    return RGBColor.from_string(value)


def add_shape(slide, kind, x, y, width, height, fill, line=LINE, radius=True):
    shape = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color(fill)
    shape.line.color.rgb = color(line)
    shape.line.width = Pt(0.7)
    return shape


def add_text(
    slide,
    text,
    x,
    y,
    width,
    height,
    size=10,
    font_color=INK,
    bold=False,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.MIDDLE,
    font_name="Segoe UI",
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(width), Inches(height))
    frame = box.text_frame
    frame.clear()
    frame.margin_left = 0
    frame.margin_right = 0
    frame.margin_top = 0
    frame.margin_bottom = 0
    frame.vertical_anchor = valign
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    paragraph.space_before = 0
    paragraph.space_after = 0
    run = paragraph.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color(font_color)
    return box


def add_line(slide, x1, y1, x2, y2, line_color=LINE, width=0.7):
    line = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    line.line.color.rgb = color(line_color)
    line.line.width = Pt(width)
    return line


def add_card(slide, x, label, value, unit, foot, accent, soft, width=2.45):
    y = 1.57
    height = 1.03
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, width, height, WHITE)
    accent_bar = add_shape(slide, MSO_SHAPE.RECTANGLE, x, y + 0.07, 0.035, height - 0.14, accent, accent)
    accent_bar.line.fill.background()
    add_text(slide, label, x + 0.14, y + 0.13, width - 0.73, 0.26, 7.8, INK, True, valign=MSO_ANCHOR.TOP)
    badge_x = x + width - 0.41
    badge = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, badge_x, y + 0.13, 0.28, 0.28, soft, soft)
    badge.line.fill.background()
    add_text(slide, "●", badge_x, y + 0.13, 0.28, 0.28, 8, accent, True, PP_ALIGN.CENTER)
    add_text(slide, value, x + 0.14, y + 0.46, 0.78, 0.33, 23, INK, True)
    add_text(slide, unit, x + 0.84, y + 0.55, 0.9, 0.15, 6.7, MUTED)
    add_line(slide, x + 0.14, y + 0.86, x + 0.28, y + 0.86, accent, 1.3)
    add_text(slide, foot, x + 0.34, y + 0.79, width - 0.49, 0.16, 5.5, MUTED)


def add_panel(slide, x, y, width, height, number, title, subtitle):
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, width, height, WHITE)
    add_line(slide, x, y + 0.48, x + width, y + 0.48, "E9EEEB", 0.6)
    number_box = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x + 0.12, y + 0.13, 0.27, 0.25, GREEN_050, GREEN_050)
    number_box.line.fill.background()
    add_text(slide, number, x + 0.12, y + 0.13, 0.27, 0.25, 6.6, GREEN_800, True, PP_ALIGN.CENTER)
    add_text(slide, title, x + 0.46, y + 0.11, width - 0.6, 0.19, 9.5, INK, True)
    add_text(slide, subtitle, x + 0.46, y + 0.30, width - 0.6, 0.11, 5.8, MUTED)


def add_demand_box(slide, x, y, label, value, description, accent):
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, 2.72, 0.65, "FAFCFB", LINE)
    bar = add_shape(slide, MSO_SHAPE.RECTANGLE, x, y + 0.04, 0.03, 0.57, accent, accent)
    bar.line.fill.background()
    add_text(slide, label.upper(), x + 0.14, y + 0.10, 2.25, 0.12, 5.6, MUTED, True)
    add_text(slide, value, x + 0.14, y + 0.25, 0.7, 0.28, 19, GREEN_900, True)
    add_text(slide, "policiais", x + 0.76, y + 0.34, 0.55, 0.12, 5.3, MUTED)
    add_text(slide, description, x + 0.14, y + 0.51, 2.4, 0.08, 4.3, MUTED)


def add_summary_box(slide, x, label, value, description, accent):
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, 1.55, 4.12, 0.68, WHITE)
    bar = add_shape(slide, MSO_SHAPE.RECTANGLE, x, 1.62, 0.035, 0.54, accent, accent)
    bar.line.fill.background()
    add_text(slide, label.upper(), x + 0.16, 1.66, 2.35, 0.12, 6.0, MUTED, True)
    add_text(slide, value, x + 3.10, 1.67, 0.72, 0.24, 18, GREEN_900, True, PP_ALIGN.RIGHT)
    add_text(slide, description, x + 0.16, 1.91, 3.62, 0.12, 5.2, MUTED)


def add_ranked_bars(slide, x, y, width, rows, total, accent, row_step=0.34, label_width=0.92):
    maximum = max(value for _, value in rows)
    track_x = x + 0.28 + label_width
    track_width = width - label_width - 1.28
    for index, (label, value) in enumerate(rows, start=1):
        row_y = y + (index - 1) * row_step
        add_text(slide, f"{index:02d}", x, row_y, 0.18, 0.14, 4.8, "9AA49F", True)
        add_text(slide, label, x + 0.23, row_y, label_width, 0.14, 5.4, "526059", True)
        track = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, track_x, row_y + 0.035, track_width, 0.075, "EAF0ED", "EAF0ED")
        track.line.fill.background()
        bar_width = max(track_width * value / maximum, 0.05)
        bar = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, track_x, row_y + 0.035, bar_width, 0.075, accent, accent)
        bar.line.fill.background()
        share = f"{value / total * 100:.1f}".replace(".", ",")
        add_text(slide, f"{value} · {share}%", x + width - 0.72, row_y, 0.72, 0.14, 5.2, INK, True, PP_ALIGN.RIGHT)


def generate_pdf() -> None:
    import pythoncom
    from win32com.client import DispatchEx

    pythoncom.CoInitialize()
    powerpoint = None
    deck = None
    try:
        powerpoint = DispatchEx("PowerPoint.Application")
        powerpoint.DisplayAlerts = 0
        deck = powerpoint.Presentations.Open(str(PPTX_PATH.resolve()), ReadOnly=True, WithWindow=False)
        deck.SaveAs(str(PDF_PATH.resolve()), 32)
    finally:
        if deck is not None:
            deck.Close()
        if powerpoint is not None:
            powerpoint.Quit()
        pythoncom.CoUninitialize()


def generate_editable_powerpoint() -> None:
    presentation = Presentation()
    presentation.slide_width = Inches(13.333333)
    presentation.slide_height = Inches(7.5)
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color(BG)

    # Cabeçalho institucional
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 0.25, 0.15, 12.83, 0.78, GREEN_900, GREEN_900)
    logo_bg = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 0.25, 0.15, 2.68, 0.78, WHITE, GREEN_900)
    logo_bg.line.width = Pt(0.7)
    slide.shapes.add_picture(str(LOGO_IMAGE), Inches(0.48), Inches(0.29), width=Inches(2.22), height=Inches(0.48))
    add_text(slide, "COMANDO-GERAL", 3.12, 0.27, 2.0, 0.13, 6.8, "74D5A6", True)
    add_text(slide, "Gestão de Pessoal e Expansão", 3.12, 0.42, 5.8, 0.28, 20, WHITE, True)
    add_text(slide, "Indicadores estratégicos consolidados · 2025–2026", 3.12, 0.70, 4.5, 0.13, 7, "B7D0C6")
    add_text(slide, "21 de agosto de 2026", 10.7, 0.35, 1.5, 0.12, 5.5, "B7D0C6", False, PP_ALIGN.RIGHT)
    add_text(slide, "16:27", 11.32, 0.53, 0.9, 0.2, 13, WHITE, True, PP_ALIGN.RIGHT)
    screen = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 12.42, 0.35, 0.36, 0.36, "234F40", "537668")
    add_text(slide, "▣", 12.42, 0.35, 0.36, 0.36, 11, WHITE, False, PP_ALIGN.CENTER)

    # Faixa de contexto
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 0.25, 1.04, 12.83, 0.40, WHITE)
    ctx_bar = add_shape(slide, MSO_SHAPE.RECTANGLE, 0.38, 1.13, 0.025, 0.22, GREEN_600, GREEN_600)
    ctx_bar.line.fill.background()
    add_text(slide, "Panorama executivo integrado", 0.49, 1.10, 2.5, 0.14, 7.1, INK, True)
    add_text(slide, "Os mostradores consolidam simultaneamente as cinco bases", 0.49, 1.25, 3.3, 0.09, 4.9, MUTED)
    add_text(slide, "REFERÊNCIA", 9.65, 1.10, 0.65, 0.09, 4.4, MUTED)
    add_text(slide, "2025–2026", 9.65, 1.22, 0.8, 0.12, 6.3, INK, True)
    add_text(slide, "ARQUIVO CONSOLIDADO", 10.75, 1.10, 1.15, 0.09, 4.4, MUTED)
    add_text(slide, "21/08/2026", 10.75, 1.22, 0.75, 0.12, 6.3, INK, True)
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 11.92, 1.13, 0.95, 0.20, GREEN_050, GREEN_050)
    add_text(slide, "●  5 bases consolidadas", 11.96, 1.13, 0.86, 0.20, 4.7, GREEN_800, True, PP_ALIGN.CENTER)

    # Cards editáveis
    card_x = [0.84, 3.80, 6.76, 9.72]
    card_width = 2.78
    add_card(slide, card_x[0], "Saídas de efetivo", "547", "saídas", "252 dem. · 88 exon. · 207 aposent.", GREEN_600, "E6F4ED", card_width)
    add_card(slide, card_x[1], "RAIO — Necessidade de efetivo das bases satélites", "912", "policiais", "20 bases · 31 municípios satélite", BLUE, "E7F1F6", card_width)
    add_card(slide, card_x[2], "Déficit de efetivo — POG", "304", "policiais", "POG — Policiamento Ostensivo Geral · 22 OPM negativas", "B58E35", "F7EFD9", card_width)
    add_card(slide, card_x[3], "COPAC — Necessidade PReVio", "229", "policiais", "Complemento para 12 bases", TEAL, "E3F3F1", card_width)

    # Painel esquerdo: gráfico editável
    add_panel(slide, 0.25, 2.72, 6.73, 4.49, "01", "Demissões e exonerações por mês", "340 das 547 saídas · janeiro a agosto de 2026")
    chart_data = ChartData()
    chart_data.categories = ["Jan", "Fev", "Mar", "Abr", "Mai", "Jun", "Jul", "Ago"]
    chart_data.add_series("Demissões", (8, 2, 0, 7, 143, 76, 14, 2))
    chart_data.add_series("Exonerações", (10, 7, 7, 24, 22, 6, 12, 0))
    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(0.46),
        Inches(3.31),
        Inches(6.25),
        Inches(3.42),
        chart_data,
    )
    chart = chart_frame.chart
    chart.has_title = False
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.TOP
    chart.legend.include_in_layout = False
    chart.legend.font.name = "Segoe UI"
    chart.legend.font.size = Pt(6)
    chart.value_axis.minimum_scale = 0
    chart.value_axis.maximum_scale = 160
    chart.value_axis.major_unit = 40
    chart.value_axis.tick_label_position = XL_TICK_LABEL_POSITION.NONE
    chart.value_axis.tick_labels.font.name = "Segoe UI"
    chart.value_axis.tick_labels.font.size = Pt(5)
    chart.value_axis.tick_labels.font.color.rgb = color("9AA49F")
    chart.value_axis.major_gridlines.format.line.color.rgb = color("E9EEEB")
    chart.category_axis.tick_labels.font.name = "Segoe UI"
    chart.category_axis.tick_labels.font.size = Pt(6)
    chart.category_axis.tick_labels.font.color.rgb = color(MUTED)
    plot = chart.plots[0]
    plot.gap_width = 65
    plot.has_data_labels = True
    plot.data_labels.show_value = True
    plot.data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
    plot.data_labels.font.name = "Segoe UI"
    plot.data_labels.font.size = Pt(5.5)
    plot.series[0].format.fill.solid()
    plot.series[0].format.fill.fore_color.rgb = color(GREEN_600)
    plot.series[0].format.line.color.rgb = color(GREEN_600)
    plot.series[1].format.fill.solid()
    plot.series[1].format.fill.fore_color.rgb = color(GOLD)
    plot.series[1].format.line.color.rgb = color(GOLD)
    add_text(slide, "Esta série detalha 340 das 547 saídas; maio concentrou 165 demissões/exonerações.", 0.46, 6.89, 6.0, 0.12, 5.1, MUTED)

    # Painel central: necessidades de efetivo
    add_panel(slide, 7.10, 2.72, 3.05, 4.49, "02", "Necessidades de efetivo", "RAIO, POG e COPAC/PReVio")
    add_demand_box(slide, 7.27, 3.34, "RAIO · Bases satélites", "912", "20 bases, 31 municípios satélite e 85,9% operacional.", BLUE)
    add_demand_box(slide, 7.27, 4.08, "Policiamento Ostensivo Geral (POG)", "304", "22 das 88 OPM apresentam saldo negativo.", "B58E35")
    add_demand_box(slide, 7.27, 4.82, "COPAC · Necessidade PReVio", "229", "Complemento para 12 bases; 63,6% do padrão projetado.", TEAL)
    warning = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 7.27, 6.72, 2.72, 0.30, "FBF8EF", "EEE4C9")
    add_text(slide, "●  304 soma os saldos negativos após excluir COGEIC, CGP e marcadores não OPM.", 7.36, 6.75, 2.55, 0.22, 4.3, "806D3C")

    # Painel direito: gráfico de rosca editável
    add_panel(slide, 10.27, 2.72, 2.81, 4.49, "03", "Aposentadorias", "Impacto nas promoções requeridas")
    donut_data = ChartData()
    donut_data.categories = ["Acesso ao oficialato", "Entre postos de oficiais"]
    donut_data.add_series("Promoções", (153, 54))
    donut_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.DOUGHNUT,
        Inches(10.65),
        Inches(3.65),
        Inches(2.05),
        Inches(2.10),
        donut_data,
    )
    donut = donut_frame.chart
    donut.has_title = False
    donut.has_legend = False
    donut.plots[0].hole_size = 64
    donut_series = donut.plots[0].series[0]
    donut_series.points[0].format.fill.solid()
    donut_series.points[0].format.fill.fore_color.rgb = color(OLIVE)
    donut_series.points[0].format.line.color.rgb = color(OLIVE)
    donut_series.points[1].format.fill.solid()
    donut_series.points[1].format.fill.fore_color.rgb = color(BLUE)
    donut_series.points[1].format.line.color.rgb = color(BLUE)
    add_text(slide, "207", 11.22, 4.29, 0.92, 0.34, 21, INK, True, PP_ALIGN.CENTER)
    add_text(slide, "Aposent.", 11.32, 4.62, 0.72, 0.12, 5.5, MUTED, False, PP_ALIGN.CENTER)
    add_shape(slide, MSO_SHAPE.RECTANGLE, 10.49, 5.92, 0.07, 0.07, OLIVE, OLIVE)
    add_text(slide, "Acesso ao oficialato", 10.62, 5.87, 1.35, 0.16, 5.3, MUTED)
    add_text(slide, "153 · 73,9%", 12.10, 5.87, 0.70, 0.16, 5.3, INK, True, PP_ALIGN.RIGHT)
    add_shape(slide, MSO_SHAPE.RECTANGLE, 10.49, 6.14, 0.07, 0.07, BLUE, BLUE)
    add_text(slide, "Entre postos de oficiais", 10.62, 6.09, 1.35, 0.16, 5.3, MUTED)
    add_text(slide, "54 · 26,1%", 12.10, 6.09, 0.70, 0.16, 5.3, INK, True, PP_ALIGN.RIGHT)

    # Rodapé
    add_text(slide, "POLÍCIA MILITAR DO CEARÁ · PAINEL ESTRATÉGICO INSTITUCIONAL", 0.25, 7.31, 4.2, 0.08, 4.2, MUTED, True)
    add_text(slide, "●  Fontes consolidadas em 21/08/2026", 10.52, 7.31, 2.55, 0.08, 4.2, GREEN_800, False, PP_ALIGN.RIGHT)

    # Segundo slide: origem territorial das saídas
    origin_slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    origin_slide.background.fill.solid()
    origin_slide.background.fill.fore_color.rgb = color(BG)

    add_shape(origin_slide, MSO_SHAPE.ROUNDED_RECTANGLE, 0.25, 0.15, 12.83, 0.78, GREEN_900, GREEN_900)
    logo_bg = add_shape(origin_slide, MSO_SHAPE.ROUNDED_RECTANGLE, 0.25, 0.15, 2.68, 0.78, WHITE, GREEN_900)
    logo_bg.line.width = Pt(0.7)
    origin_slide.shapes.add_picture(str(LOGO_IMAGE), Inches(0.48), Inches(0.29), width=Inches(2.22), height=Inches(0.48))
    add_text(origin_slide, "GESTÃO DE PESSOAL", 3.12, 0.27, 2.0, 0.13, 6.8, "74D5A6", True)
    add_text(origin_slide, "De onde estão saindo os militares", 3.12, 0.42, 6.2, 0.28, 19, WHITE, True)
    add_text(origin_slide, "Origem por OPM e município · dados agregados, sem nomes ou matrículas", 3.12, 0.70, 5.5, 0.13, 7, "B7D0C6")
    add_text(origin_slide, "PERÍODO DA FONTE", 10.70, 0.33, 1.30, 0.10, 4.6, "B7D0C6", False, PP_ALIGN.RIGHT)
    add_text(origin_slide, "01/01 a 10/08/2026", 10.42, 0.52, 1.58, 0.15, 8.2, WHITE, True, PP_ALIGN.RIGHT)
    add_shape(origin_slide, MSO_SHAPE.ROUNDED_RECTANGLE, 12.42, 0.35, 0.36, 0.36, "234F40", "537668")
    add_text(origin_slide, "04", 12.42, 0.35, 0.36, 0.36, 8, WHITE, True, PP_ALIGN.CENTER)

    add_shape(origin_slide, MSO_SHAPE.ROUNDED_RECTANGLE, 0.25, 1.04, 12.83, 0.38, WHITE)
    ctx_bar = add_shape(origin_slide, MSO_SHAPE.RECTANGLE, 0.38, 1.12, 0.025, 0.22, GREEN_600, GREEN_600)
    ctx_bar.line.fill.background()
    add_text(origin_slide, "Leitura territorial das saídas", 0.49, 1.09, 2.3, 0.14, 7.1, INK, True)
    add_text(origin_slide, "O ranking cobre demissões e exonerações individualizadas; aposentadorias não possuem OPM na base atual.", 0.49, 1.24, 6.4, 0.09, 4.9, MUTED)
    add_shape(origin_slide, MSO_SHAPE.ROUNDED_RECTANGLE, 11.46, 1.12, 1.38, 0.20, GREEN_050, GREEN_050)
    add_text(origin_slide, "326 registros únicos", 11.50, 1.12, 1.30, 0.20, 5.2, GREEN_800, True, PP_ALIGN.CENTER)

    add_summary_box(origin_slide, 0.25, "Registros com origem", "326", "Demissões e exonerações individualizadas", GREEN_600)
    add_summary_box(origin_slide, 4.60, "Unidades identificadas", "65", "OPMs consolidadas pela unidade principal", BLUE)
    add_summary_box(origin_slide, 8.95, "Municípios identificados", "53", "Fortaleza concentra 153 registros", GOLD)

    add_panel(origin_slide, 0.25, 2.39, 6.28, 4.72, "04A", "OPMs com mais saídas", "Companhias agrupadas pela unidade principal · Top 10")
    opm_rows = [
        ("CPRAIO", 17), ("12º BPM", 17), ("17º BPM", 14), ("18º BPM", 14), ("6º BPM", 14),
        ("BPTUR", 13), ("COPAC", 13), ("19º BPM", 12), ("20º BPM", 11), ("24º BPM", 10),
    ]
    add_ranked_bars(origin_slide, 0.48, 3.10, 5.80, opm_rows, 326, GREEN_600, row_step=0.325, label_width=0.86)
    add_text(origin_slide, "CPRAIO e 12º BPM lideram, com 17 registros cada (5,2% do recorte).", 0.48, 6.55, 5.75, 0.14, 5.2, MUTED)

    add_panel(origin_slide, 6.68, 2.39, 6.40, 4.72, "04B", "Municípios com mais saídas", "Município informado no registro · Top 8")
    city_rows = [
        ("Fortaleza", 153), ("Caucaia", 27), ("Maracanaú", 12), ("Maranguape", 8),
        ("Juazeiro do Norte", 8), ("Quixadá", 7), ("Eusébio", 6), ("Sobral", 6),
    ]
    add_ranked_bars(origin_slide, 6.91, 3.10, 5.91, city_rows, 326, BLUE, row_step=0.39, label_width=1.22)
    add_text(origin_slide, "Fortaleza reúne 153 registros, equivalentes a 46,9% do recorte individualizado.", 6.91, 6.55, 5.75, 0.14, 5.2, MUTED)

    note = add_shape(origin_slide, MSO_SHAPE.ROUNDED_RECTANGLE, 0.48, 6.78, 12.34, 0.22, "FBF8EF", "EEE4C9")
    note.line.width = Pt(0.5)
    add_text(origin_slide, "Cobertura metodológica: 326 registros únicos. O total executivo de 340 segue a base mensal consolidada; as 207 aposentadorias não possuem origem territorial informada.", 0.60, 6.81, 12.08, 0.14, 4.8, "806D3C", False, PP_ALIGN.CENTER)

    add_text(origin_slide, "POLÍCIA MILITAR DO CEARÁ · PAINEL ESTRATÉGICO INSTITUCIONAL", 0.25, 7.31, 4.2, 0.08, 4.2, MUTED, True)
    add_text(origin_slide, "●  Dados agregados · sem exposição de informações pessoais", 9.45, 7.31, 3.63, 0.08, 4.2, GREEN_800, False, PP_ALIGN.RIGHT)

    presentation.core_properties.title = "Painel de Gestão de Pessoal e Expansão - PMCE"
    presentation.core_properties.subject = "Indicadores estratégicos consolidados de 2025-2026"
    presentation.core_properties.author = "Polícia Militar do Ceará"
    presentation.core_properties.comments = "Dois slides com elementos, textos e gráficos editáveis em formato 16:9."
    presentation.save(PPTX_PATH)


def main() -> None:
    for required in (LOGO_IMAGE,):
        if not required.exists():
            raise FileNotFoundError(f"Arquivo de origem não encontrado: {required}")
    EXPORTS.mkdir(exist_ok=True)
    generate_editable_powerpoint()
    generate_pdf()
    print(f"PDF: {PDF_PATH}")
    print(f"PowerPoint editável: {PPTX_PATH}")


if __name__ == "__main__":
    main()
